import json
import random
from pptx import Presentation
import requests
import re
from collections import Counter

# Define the ABCD ID list
abcd_id_list = [
    711, 733, 523, 542, 731, 763, 669, 739, 690, 742,
    724, 112, 668, 714, 671, 32, 722, 434, 713, 751,
    760, 102, 654, 695, 762, 470, 33, 432, 429, 265,
    673, 111, 588, 734, 692, 477, 518, 720, 445, 437,
    752, 206, 670, 544, 491, 737, 685, 492, 53, 676,
    575, 469, 183, 717, 665, 715, 410, 425, 430, 586,
    769, 235, 462, 741, 524, 718, 758, 465, 745, 520,
    431, 721, 171, 484, 759, 672, 433, 749, 502, 756,
    747, 757, 26, 765, 667
]


# Function to fetch text data for a given ABCD ID
def fetch_text_data(abcd_id):
    response = requests.get(f'https://abcd2.projectabcd.com/api/getinfo.php?id={abcd_id}', headers={"User-Agent": "XY"})
    json_response = json.loads(response.text)
    if 'data' in json_response:
        data = json_response['data']
        text_data = ""
        for key in ['description', 'did_you_know']:
            if key in data and data[key]:
                text_data += data[key] + " "
        return text_data.strip()
    else:
        return None


# Function to interactively modify the word list
def modify_word_list(word_list):
    print("Current word list:")
    for i, word in enumerate(word_list):
        print(f"{i + 1}. {word}")
    while True:
        action = input("Choose an action: (A)dd, (D)elete, (M)odify, (F)inish, (E)nd: ").lower()
        if action == 'a':
            new_word = input("Enter the new word: ")
            word_list.append(new_word)
        elif action == 'd':
            index = int(input("Enter the index of the word to delete: ")) - 1
            if 0 <= index < len(word_list):
                del word_list[index]
        elif action == 'm':
            index = int(input("Enter the index of the word to modify: ")) - 1
            if 0 <= index < len(word_list):
                new_word = input("Enter the new word: ")
                word_list[index] = new_word
        elif action == 'e':
            return word_list, False
        elif action == 'f':
            return word_list, True
        else:
            print("Invalid action. Please choose again.")


# Function to generate the word search grid
def generate_word_search_grid(words_for_puzzle):
    puzzle_size = 10  # Size of the puzzle grid
    puzzle_grid = [[' ' for _ in range(puzzle_size)] for _ in range(puzzle_size)]
    for word in words_for_puzzle:
        direction = random.choice(['horizontal', 'vertical', 'diagonal'])
        if direction == 'horizontal':
            pass  # Logic to place word horizontally
        elif direction == 'vertical':
            pass  # Logic to place word vertically
        else:
            pass  # Logic to place word diagonally
    return puzzle_grid


# Function to create PowerPoint presentation
def create_powerpoint_presentation(puzzle, solution):
    prs = Presentation()
    # Add puzzle slide
    puzzle_slide = prs.slides.add_slide(prs.slide_layouts[5])
    puzzle_title = puzzle_slide.shapes.title
    puzzle_title.text = "Word Search Puzzle"
    # Check if placeholder exists
    if len(puzzle_slide.placeholders) > 1:
        puzzle_text = puzzle_slide.placeholders[1].text_frame
        puzzle_text.text = puzzle
    else:
        print("No placeholder found on puzzle slide.")

    # Add solution slide
    solution_slide = prs.slides.add_slide(prs.slide_layouts[5])
    solution_title = solution_slide.shapes.title
    solution_title.text = "Solution"
    # Check if placeholder exists
    if len(solution_slide.placeholders) > 1:
        solution_text = solution_slide.placeholders[1].text_frame
        solution_text.text = solution
    else:
        print("No placeholder found on solution slide.")

    prs.save("word_search_puzzle.pptx")


if __name__ == "__main__":
    # Fetch text data for each ABCD ID
    text_data_dict = {}
    for abcd_id in abcd_id_list:
        text_data = fetch_text_data(abcd_id)
        text_data_dict[abcd_id] = text_data

    # Let the user modify the word list
    modified_word_lists = {}
    for abcd_id, text_data in text_data_dict.items():
        if text_data:
            words = re.findall(r'\b\w+\b', text_data.lower())  # Tokenize words
            word_freq = Counter(words)
            initial_word_list = list(word_freq.keys())[:10]  # Select top 10 words initially
            modify_next_id = True
            while modify_next_id:
                modified_word_list, modify_next_id = modify_word_list(initial_word_list)
            modified_word_lists[abcd_id] = modified_word_list

            if not modify_next_id:
                break  # Stop modifying IDs if the user chooses to end

    # Generate word search puzzle for the final word list
    words_for_puzzle = []
    for word_list in modified_word_lists.values():
        words_for_puzzle.extend(word_list)

    puzzle_grid = generate_word_search_grid(words_for_puzzle)

    # Convert puzzle grid to string format for presentation
    puzzle_str = '\n'.join([' '.join(row) for row in puzzle_grid])
    solution_str = "Solution string"  # Replace this with the actual solution string

    # Create PowerPoint presentation
    create_powerpoint_presentation(puzzle_str, solution_str)
