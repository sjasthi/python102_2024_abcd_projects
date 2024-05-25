import random
import string
import requests
import json
import re
import sheroes_base_lib as sheroes
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Function to fetch text data for a given ABCD ID
def fetch_text_data(abcd_id):
    response = requests.get(f'https://abcd2.projectabcd.com/api/getinfo.php?id={abcd_id}', headers={"User-Agent": "XY"})
    json_response = json.loads(response.text)
    if 'data' in json_response:
        data = json_response['data']
        text_data = ""
        for key in ['description', 'did_you_know']:
            if key in data and data[key]:
                text_data += data[key] + " "
        return text_data.strip()
    else:
        return None

# Function to identify words from text
def identify_words(text):
    words = re.findall(r'\b\w+\b', text.lower())
    return words

# Function to allow user to add, modify, delete, or finish words in intervals of 10
def edit_word_list(words):
    i = 0
    while i < len(words):
        sublist = words[i:i + 10]
        while True:
            print("Current word list segment:")
            for index, word in enumerate(sublist, start=1):
                print(f"{index}. {word}")

            action = input(
                "Do you want to (A)dd, (M)odify, (D)elete, (F)inish this segment, or (E)nd editing? ").strip().upper()
            if action == "A":
                new_word = input("Enter word to add: ").strip().lower()
                sublist.append(new_word)
            elif action == "M":
                index = int(input("Enter index of word to modify: "))
                if 1 <= index <= len(sublist):
                    new_word = input("Enter new word: ").strip().lower()
                    sublist[index - 1] = new_word
                else:
                    print("Invalid index.")
            elif action == "D":
                index = int(input("Enter index of word to delete: "))
                if 1 <= index <= len(sublist):
                    sublist.pop(index - 1)
                else:
                    print("Invalid index.")
            elif action == "F":
                break
            elif action == "E":
                return True
            else:
                print("Invalid action. Please choose A, M, D, F, or E.")

            print("Updated word list segment:")
            for index, word in enumerate(sublist, start=1):
                print(f"{index}. {word}")

        words[i:i + 10] = sublist
        i += 10
    return False

# Function to create a word search puzzle grid
def create_puzzle_grid(size):
    return [[' ' for _ in range(size)] for _ in range(size)]

# Function to add words to the puzzle grid
def add_word_to_grid(word, grid):
    directions = [(1, 0), (0, 1), (-1, 0), (0, -1), (1, 1), (-1, 1), (1, -1), (-1, -1)]
    direction = random.choice(directions)
    word_length = len(word)
    grid_size = len(grid)

    # Generate random starting position
    row = random.randint(0, grid_size - 1)
    col = random.randint(0, grid_size - 1)

    # Check if the word can fit in the grid from the starting position in the chosen direction
    while not (0 <= row + (word_length - 1) * direction[0] < grid_size and 0 <= col + (word_length - 1) * direction[1] < grid_size):
        direction = random.choice(directions)
        row = random.randint(0, grid_size - 1)
        col = random.randint(0, grid_size - 1)

    # Place the word in the grid
    for i in range(word_length):
        grid[row + i * direction[0]][col + i * direction[1]] = word[i]

# Function to fill the remaining grid cells with random letters
def fill_grid_random(grid):
    for i in range(len(grid)):
        for j in range(len(grid[i])):
            if grid[i][j] == ' ':
                grid[i][j] = random.choice(string.ascii_uppercase)

# Function to generate a word search puzzle
def generate_word_search(words, grid_size):
    grid = create_puzzle_grid(grid_size)
    for word in words:
        add_word_to_grid(word.upper(), grid)
    solution_grid = [row[:] for row in grid]  # Copy grid for solution
    fill_grid_random(grid)
    return grid, solution_grid

# Function to print the puzzle grid
def print_grid(grid):
    for row in grid:
        print(' '.join(row))

# Function to add a slide with the word search puzzle to the presentation
def add_puzzle_slide(prs, grid, abcd_id):
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = f"Label #{abcd_id}:"

    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(6))
    text_frame = textbox.text_frame
    p = text_frame.add_paragraph()
    for row in grid:
        p.text += ' '.join(row) + '\n'
    p.font.name = 'Courier New'
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.LEFT

# Function to add a slide with the solution to the word search puzzle to the presentation
def add_solution_slide(prs, solution_grid, abcd_id, words):
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = f"Label #{abcd_id} Solution:"

    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(6))
    text_frame = textbox.text_frame
    p = text_frame.add_paragraph()

    for word in words:
        for row_index, row in enumerate(solution_grid):
            for col_index, letter in enumerate(row):
                if letter.lower() == word[0].lower():  # Check if the letter matches the first letter of the word
                    # Check if the word fits horizontally
                    if col_index + len(word) <= len(row) and ''.join(row[col_index:col_index+len(word)]).lower() == word.lower():
                        for i in range(col_index, col_index + len(word)):
                            solution_grid[row_index][i] = word[i - col_index]
                    # Check if the word fits vertically
                    elif row_index + len(word) <= len(solution_grid) and ''.join(solution_grid[i][col_index] for i in range(row_index, row_index+len(word))).lower() == word.lower():
                        for i in range(row_index, row_index + len(word)):
                            solution_grid[i][col_index] = word[i - row_index]
                    # Check if the word fits diagonally (down-right)
                    elif col_index + len(word) <= len(row) and row_index + len(word) <= len(solution_grid) and ''.join(solution_grid[row_index+i][col_index+i] for i in range(len(word))).lower() == word.lower():
                        for i in range(len(word)):
                            solution_grid[row_index+i][col_index+i] = word[i]

    for row in solution_grid:
        p.text += ' '.join(row) + '\n'
    p.font.name = 'Courier New'
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.LEFT

# Main function
def main():
    fetched_data_list = sheroes.fetch_data_from_api()

    word_lists = {}

    for data in fetched_data_list:
        text_data = data.get('description', '') + ' ' + data.get('did_you_know', '')
        if text_data:
            words = identify_words(text_data)
            word_lists[data['id']] = words

    prs = Presentation()

    for abcd_id, words in word_lists.items():
        print(f"Editing word list for ABCD ID {abcd_id}:")
        while True:
            if edit_word_list(words):
                break

        grid_size = max(len(word) for word in words) * 2
        puzzle, solution = generate_word_search(words, grid_size)

        print(f"Word Search Puzzle for ABCD ID {abcd_id}:")
        print_grid(puzzle)
        print()

        add_puzzle_slide(prs, puzzle, abcd_id)
        add_solution_slide(prs, solution, abcd_id, words)

    prs.save('word_search_puzzles.pptx')
    print("Word search puzzles saved to word_search_puzzles.pptx")

# Execute the main function
if __name__ == "__main__":
    main()


