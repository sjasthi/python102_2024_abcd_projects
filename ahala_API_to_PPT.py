# Importing the necessary Python libraries
import json
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
import pptx.util

# Helper function
def find_content_placeholder_index(Layout):
    for idx, shape in enumerate(Layout.placeholders):
        if shape.placeholder_format.idx == 1:
            return idx
    return None

# Defining the function
def get_sheroes_ppt(abcd_id_list):

    # Using prs to create a presentation and slide with specific parameters
    prs = Presentation()
    prs.slide_width = pptx.util.Inches(11)
    prs.slide_height = pptx.util.Inches(9)

    # Getting the json response (did you know, description, name) from the Sheroes API ids
    for id_number in abcd_id_list:
        response = requests.get(f'https://abcd2.projectabcd.com/api/getinfo.php?id={id_number}', headers={"User-Agent": "XY"})
        json_response = json.loads(response.text)

        desc = json_response['data']['description']
        dyk = json_response['data']['did_you_know']
        name = json_response['data']['name']
        text = 'Description: ' + desc + ' Did You Know: ' + dyk

        # Specifying the layout - we have used layout 3 here
        Layout = prs.slide_layouts[3]
        slide = prs.slides.add_slide(Layout)

        # Adding the name as the title of the slide
        slide.shapes.title.text = name
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(40)
        slide.placeholders[1].text = text

        # Minimizing the font size of the text in the placeholder 1
        content_placeholder_idx = find_content_placeholder_index(Layout)
        if content_placeholder_idx is not None:
            content_box = slide.placeholders[content_placeholder_idx]
            text_frame = content_box.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(13)

        # I have locally imported only 5 images
        # So, the rest will have a generic picture
        if id_number == 26:
            img_path = 'C:/Users/ugant/OneDrive/Desktop/Python102/26.png'

        elif id_number == 27:
            img_path = 'C:/Users/ugant/OneDrive/Desktop/Python102/27.png'

        elif id_number == 28:
            img_path = 'C:/Users/ugant/OneDrive/Desktop/Python102/28.png'

        elif id_number == 29:
            img_path = 'C:/Users/ugant/OneDrive/Desktop/Python102/29.png'

        elif id_number == 30:
            img_path = 'C:/Users/ugant/OneDrive/Desktop/Python102/30.png'

        else:
            img_path = 'C:/Users/ugant/OneDrive/Desktop/Python102/sheroes_cover.png'

        # Parameters for the image
        left_inch = Inches(5.00)
        top_inch = Inches(1.75)
        width_inch = Inches(5.5)
        height_inch = Inches(6.25)

        slide.shapes.add_picture(str(img_path), left_inch, top_inch, width_inch, height_inch)

    # Saving the presentation
    prs.save("sheroes_api_to_ppt.pptx")
    print("done")

# Giving the id_list and calling the function
id_list = [26, 27, 28, 29, 30, 31, 32, 33, 39, 50, 52, 53, 101, 102, 110, 111, 112, 114, 115, 116, 117, 119, 151,
           171, 183, 188, 196, 201, 206, 235, 265, 275, 276, 306, 313, 314, 316, 317, 318, 319, 320, 321, 322, 323,
           324, 325, 326, 327, 328, 329, 330, 356, 401, 405, 406, 407, 409, 410, 411, 412, 413, 414, 415, 418, 419,
           420, 422, 423, 424, 425, 426, 427, 428, 429, 430, 431, 432, 433, 434, 435, 436, 437, 438, 439, 440, 441,
           442, 443, 444, 445, 462, 463, 464, 465, 466, 467, 468, 469, 470, 471, 472, 473, 474, 475, 476, 477, 478,
           483, 484, 491, 492, 493, 502, 506, 518, 520, 523, 524, 525, 534, 538, 542, 544, 547, 548, 549, 560, 568,
           574, 575, 577, 578, 581, 582, 583, 586, 588, 601, 605, 611, 618, 620, 626, 631, 632, 654, 655, 658, 659,
           660, 662, 663, 664, 665, 666, 667, 668, 669, 670, 671, 672, 673, 674, 676, 677, 678, 679, 680, 681, 682,
           683, 684, 685, 686, 687, 688, 689, 690, 691, 692, 693, 694, 695, 696, 697, 698, 699, 700, 701, 702, 703,
           704, 706, 709, 710, 711, 712, 713, 714, 715, 716, 717, 718, 719, 720, 721, 722, 723, 724, 725, 726, 728,
           729, 730, 731, 732, 733, 734, 735, 736, 737, 738, 739, 740, 741, 742, 743, 744, 745, 746, 747, 748, 749,
           750, 751, 752, 753, 754, 755, 756, 757, 758, 759, 760, 761, 762, 763, 764, 765, 766, 767, 768, 769, 770,
           771, 772]
get_sheroes_ppt(id_list)
